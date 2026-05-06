"""
SiloForge Presentation Builder
Generates SiloForge_Presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (matches SiloForge design system) ──────────────────
DARK_GREEN  = RGBColor(0x1b, 0x43, 0x32)
MED_GREEN   = RGBColor(0x2d, 0x6a, 0x4f)
LIGHT_GREEN = RGBColor(0x74, 0xc6, 0x9d)
PALE_GREEN  = RGBColor(0xd8, 0xf3, 0xdc)
PURPLE_DARK = RGBColor(0x4a, 0x10, 0x60)
PURPLE_MED  = RGBColor(0x7b, 0x2d, 0x8b)
PURPLE_PALE = RGBColor(0xf3, 0xe8, 0xfd)
AMBER       = RGBColor(0xe9, 0xa3, 0x19)
AMBER_BG    = RGBColor(0xfe, 0xf3, 0xcd)
RED         = RGBColor(0xd6, 0x28, 0x39)
RED_BG      = RGBColor(0xfd, 0xe8, 0xea)
BLUE        = RGBColor(0x48, 0x95, 0xef)
BLUE_BG     = RGBColor(0xe8, 0xf4, 0xfd)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG    = RGBColor(0xf0, 0xf4, 0xf1)
LIGHT_BG2   = RGBColor(0xf5, 0xf0, 0xf7)
DARK_TEXT   = RGBColor(0x1a, 0x1a, 0x1a)
GRAY        = RGBColor(0x6c, 0x75, 0x7d)
LIGHT_GRAY  = RGBColor(0xe8, 0xe8, 0xe8)

# ── Paths ──────────────────────────────────────────────────────────────
BASE  = r"C:\Users\FARSH\OneDrive\Desktop\SiloForge"
IMGS  = os.path.join(BASE, "images")
LAM   = os.path.join(IMGS, "lamuad_1")

def P(name):
    return os.path.join(IMGS, name)

def L(name):
    return os.path.join(LAM, name)

# ── Helpers ─────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def oval(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box

def mtxt(slide, lines, l, t, w, h,
         size=13, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT):
    """Multi-line text box; each item in lines is str or (str, size, bold, color)."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        if isinstance(line, tuple):
            run.text    = line[0]
            run.font.size  = Pt(line[1])
            run.font.bold  = line[2] if len(line) > 2 else bold
            run.font.color.rgb = line[3] if len(line) > 3 else color
        else:
            run.text = line
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = color

def img(slide, path, l, t, w=None, h=None):
    if not os.path.exists(path):
        print(f"  WARNING: image not found – {path}")
        return
    kw = {}
    if w: kw["width"]  = Inches(w)
    if h: kw["height"] = Inches(h)
    slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)

# ── Build presentation ──────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

rect(s,  0,    0,   13.33, 7.5,  DARK_GREEN)
rect(s,  0,    0,   0.55,  7.5,  MED_GREEN)   # accent left bar
oval(s, 10.0, -1.8,  5.5,  5.5,  MED_GREEN)   # top-right deco circle
oval(s, 11.2,  5.0,  3.2,  3.2,  RGBColor(0x40,0x91,0x6c))  # bottom-right deco

txt(s, "SiloForge",
    0.9, 1.1, 11, 1.5, size=68, bold=True, color=WHITE)
txt(s, "Agri-Business Intelligence for Nigerian Poultry Farms",
    0.9, 2.75, 10, 0.7, size=21, color=LIGHT_GREEN)
rect(s, 0.9, 3.58, 5.5, 0.045, LIGHT_GREEN)   # separator

txt(s, "Kamorudeen Fasasi",
    0.9, 3.8, 9, 0.5, size=16, bold=True, color=WHITE)
txt(s, "Agricultural Engineer  ·  Data Engineer  ·  Poultry Farmer",
    0.9, 4.35, 10.5, 0.4, size=13, color=LIGHT_GREEN, italic=True)
txt(s, "May 2026", 0.9, 5.0, 4, 0.35, size=12, color=GRAY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INDUSTRY FACTS
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.05, DARK_GREEN)

txt(s, "Nigeria's Poultry Industry: Big Market. Untapped Potential.",
    0.35, 0.14, 12.6, 0.78, size=22, bold=True, color=WHITE)

stats = [
    ("~$2.5B",  "Estimated annual\nvalue of Nigerian\npoultry industry"),
    ("~180M",   "Birds raised\nacross Nigeria\neach year"),
    ("60–70%",  "Of total farm cost\nis feed — highest\nsingle expense"),
    (">500k MT","Poultry products\nimported yearly\ndue to supply gap"),
    ("<5%",     "SME farms with\nconsistent\ndigital records"),
]
for i, (val, lbl) in enumerate(stats):
    x = 0.35 + i * 2.56
    rect(s, x,      1.25, 2.36, 2.1,  WHITE)
    rect(s, x,      1.25, 2.36, 0.07, MED_GREEN)
    txt(s, val,  x+0.1, 1.38, 2.2, 0.75, size=28, bold=True,
        color=MED_GREEN, align=PP_ALIGN.CENTER)
    txt(s, lbl,  x+0.1, 2.15, 2.2, 0.9,  size=10.5,
        color=DARK_TEXT, align=PP_ALIGN.CENTER)

txt(s, "The gap is not just a farming problem — it's a data problem.",
    0.5, 3.6, 12.3, 0.55, size=19, bold=True,
    color=DARK_GREEN, align=PP_ALIGN.CENTER)
rect(s, 2.8, 3.56, 7.7, 0.05, LIGHT_GREEN)
rect(s, 2.8, 4.18, 7.7, 0.05, LIGHT_GREEN)

bullets = [
    "Most farms cannot tell you their production rate from last week.",
    "Feed is the biggest cost — and variance goes undetected for weeks.",
    "No records = no access to grants, loans, or institutional buyers.",
]
for i, b in enumerate(bullets):
    y = 4.35 + i * 0.65
    rect(s, 0.5, y+0.08, 0.12, 0.3, AMBER)
    txt(s, b, 0.78, y, 11.5, 0.55, size=14, color=DARK_TEXT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE DATA ANALOGY
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, WHITE)
rect(s, 0, 0, 13.33, 1.05, DARK_GREEN)

txt(s, "Every Modern Industry Runs on Data. Why Not Your Farm?",
    0.35, 0.14, 12.6, 0.78, size=22, bold=True, color=WHITE)

cols_data = [
    ("🏦", "Banking",  RGBColor(0x1a,0x5f,0xa8), BLUE_BG,
     [" Your bank tracks every naira in real time,",
      " flags unusual transactions within seconds,",
      " and personalises your credit limit.",
      "",
      " They don't guess your balance.",
      " They don't wait until year-end",
      " to know if they made a profit."]),
    ("📡", "Telecom",  PURPLE_MED, PURPLE_PALE,
     [" MTN and Airtel know exactly which",
      " towers are underperforming, which",
      " customers are about to churn.",
      "",
      " They predict demand, allocate resources,",
      " and act before problems happen.",
      " All in real time."]),
    ("🐔", "Your Farm", DARK_GREEN, PALE_GREEN,
     [" Right now, most decisions are made",
      " from memory, paper records,",
      " or a worker's verbal report.",
      "",
      " Production rate, feed waste, mortality",
      " trends, cost per egg —",
      " most farms don't track these numbers."]),
]
for i, (icon, title, hdr_col, bg_col, lines) in enumerate(cols_data):
    x = 0.35 + i * 4.28
    rect(s, x, 1.2, 3.98, 5.7, bg_col)
    rect(s, x, 1.2, 3.98, 0.62, hdr_col)
    txt(s, f"{icon}  {title}", x+0.18, 1.26, 3.7, 0.5,
        size=17, bold=True, color=WHITE)
    for j, line in enumerate(lines):
        txt(s, line, x+0.15, 1.98 + j*0.53, 3.72, 0.5, size=12.5, color=DARK_TEXT)

rect(s, 0, 6.82, 13.33, 0.68, MED_GREEN)
txt(s, ("The smartphone in your pocket is already the analytics platform."
        "  SiloForge makes it work for your farm."),
    0.4, 6.88, 12.5, 0.55, size=13, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE "BEFORE" PICTURE
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.05, RED)

txt(s, "This is How a 3,000-Bird Farm Tracks Its Operation Today",
    0.35, 0.14, 12.6, 0.78, size=22, bold=True, color=WHITE)

photo_data = [
    (L("c.jpeg"), "Pre-printed production chart\nfilled in by hand daily"),
    (L("a.jpeg"), "Loose paper — feed &\nmedication notes"),
    (L("b.jpeg"), "Daily observations on\nscrap paper"),
]
for i, (path, caption) in enumerate(photo_data):
    x = 0.35 + i * 4.3
    rect(s, x, 1.15, 3.95, 4.85, WHITE)
    img(s,  path, x+0.1, 1.22, 3.75, 3.95)
    txt(s, caption, x, 5.15, 3.95, 0.7,
        size=10.5, color=GRAY, align=PP_ALIGN.CENTER)

rect(s, 0, 6.02, 13.33, 1.48, RED_BG)
rect(s, 0, 6.02, 0.55, 1.48, RED)
txt(s, ("No trend analysis.   No real-time alerts.   No aggregation across pens.   "
        "No cost-per-bird.   No historical baseline."),
    0.7, 6.12, 12.3, 0.48, size=13.5, bold=True,
    color=RED, align=PP_ALIGN.CENTER)
txt(s, ("When a disease outbreak hits — there is nothing to compare against. "
        "When an investor asks for records — there is nothing to show."),
    0.7, 6.65, 12.3, 0.6, size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — COST OF NO DATA
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0,   0, 13.33, 7.5, WHITE)
rect(s, 0,   0,  2.9,  7.5, DARK_GREEN)   # sidebar

txt(s, "What Poor\nRecord-\nKeeping\nCosts You",
    0.12, 0.75, 2.6, 3.2, size=22, bold=True, color=WHITE)
rect(s, 0.28, 4.05, 2.0, 0.055, LIGHT_GREEN)
txt(s, "Every day\nwithout data\nis a day of\nguesswork.",
    0.18, 4.22, 2.6, 1.6, size=13, color=LIGHT_GREEN, italic=True)

costs = [
    (AMBER, "Feed leakage",
     "Feed = 60–70% of your cost. Under-feeding and over-feeding go undetected for weeks."),
    (RED,   "Performance blind spot",
     "You can't pinpoint which pen or cage is pulling down your overall production rate."),
    (RED,   "Disease response delay",
     "No health baseline means outbreaks escalate before you notice a pattern."),
    (BLUE,  "Zero access to finance",
     "Banks, investors, and grant bodies ask for production records — most farms have none."),
    (AMBER, "Gut-feel decisions",
     "Restocking, culling, supply allocation — all based on memory, not evidence."),
]
for i, (col, heading, detail) in enumerate(costs):
    y = 0.45 + i * 1.2
    rect(s, 3.1, y, 0.07, 0.85, col)
    txt(s, heading, 3.3, y,       9.6, 0.36, size=14.5, bold=True, color=DARK_GREEN)
    txt(s, detail,  3.3, y+0.38,  9.6, 0.55, size=12,   color=DARK_TEXT)

rect(s, 3.1, 6.4, 9.7, 0.06, LIGHT_GREEN)
txt(s, ("The farms that will scale are the ones that make data-driven decisions "
        "— at every flock, every week."),
    3.1, 6.55, 9.7, 0.7, size=13, bold=True,
    color=MED_GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — INTRODUCING SILOFORGE
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, DARK_GREEN)
oval(s,  9.6, -1.5, 5.5, 5.5, MED_GREEN)
oval(s, 11.1,  5.1, 3.2, 3.2, RGBColor(0x40,0x91,0x6c))
oval(s, -1.0,  5.5, 2.8, 2.8, RGBColor(0x40,0x91,0x6c))

txt(s, "Introducing SiloForge",
    0.5, 0.32, 10, 0.85, size=36, bold=True, color=WHITE)
txt(s, ("A free, mobile-first farm management suite — "
        "built by a farmer who lives this problem."),
    0.5, 1.2, 10.5, 0.55, size=15.5, color=LIGHT_GREEN, italic=True)
rect(s, 0.5, 1.88, 7.5, 0.05, LIGHT_GREEN)

# LayerTrack card
rect(s, 0.5, 2.08, 5.95, 4.75, MED_GREEN)
rect(s, 0.5, 2.08, 5.95, 0.58, RGBColor(0x40,0x91,0x6c))
txt(s, "🌿  LayerTrack", 0.68, 2.13, 5.7, 0.5, size=18, bold=True, color=WHITE)
layer_feat = [
    "• Daily egg collection — 3 rounds per day",
    "• Production rate & crack rate with live alerts",
    "• Cell-level view — top vs underperforming cells",
    "• Flock stage tracking (Pre-Lay → Peak → Decline)",
    "• Feed log with daily variance vs requirement",
    "• Health log — water, droppings, medications",
    "• Sales, receivables and full P&L",
    "• 7-day trend reports",
]
for i, feat in enumerate(layer_feat):
    txt(s, feat, 0.72, 2.8+i*0.49, 5.65, 0.46, size=11.5, color=PALE_GREEN)

# BroodTrack card
rect(s, 6.85, 2.08, 5.95, 4.75, PURPLE_DARK)
rect(s, 6.85, 2.08, 5.95, 0.58, PURPLE_MED)
txt(s, "🐣  BroodTrack", 7.03, 2.13, 5.7, 0.5, size=18, bold=True, color=WHITE)
brood_feat = [
    "• Batch lifecycle — DOC arrival to POL sale",
    "• Vaccination schedule with overdue alerts",
    "• Weight vs breed benchmark (ISA, Lohmann, Noiler…)",
    "• Temperature monitoring during brooding phase",
    "• Feed phase management & FCR tracking",
    "• Survival rate per batch",
    "• Per-batch profit & loss",
    "• Breed benchmarks built in",
]
for i, feat in enumerate(brood_feat):
    txt(s, feat, 7.05, 2.8+i*0.49, 5.65, 0.46, size=11.5, color=PURPLE_PALE)

rect(s, 0, 7.05, 13.33, 0.45, RGBColor(0x12, 0x30, 0x22))
txt(s, ("Works on any smartphone  ·  No app download  ·  "
        "Works offline  ·  Real-time sync across devices"),
    0.4, 7.09, 12.5, 0.38, size=12, bold=True,
    color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — LAYERTRACK IN ACTION
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.0,  DARK_GREEN)

txt(s, "LayerTrack: Every Egg Accounted For",
    0.35, 0.14, 12.6, 0.73, size=24, bold=True, color=WHITE)

callouts = [
    (LIGHT_GREEN, "Production rate per pen and per cell",
     "Know instantly which cage is underperforming — not at month-end."),
    (LIGHT_GREEN, "Three daily collection rounds",
     "Morning, afternoon, evening — each round tracked with totals and broken eggs."),
    (AMBER,       "Crack rate alert",
     "Automatic red flag when crack rate exceeds 3% — signals feed or handling issue."),
    (AMBER,       "Feed variance vs daily requirement",
     "Green / amber / red bars show immediately if birds are being underfed."),
    (LIGHT_GREEN, "Finance: cost, sales, margin",
     "₦39,500 sales vs ₦67,800 costs — know your margin in real time, not year-end."),
    (BLUE,        "7-day and trend reports",
     "Spot production decline weeks before it becomes a problem."),
]
for i, (col, heading, detail) in enumerate(callouts):
    y = 1.1 + i * 1.02
    rect(s, 0.35, y+0.08, 0.06, 0.7, col)
    txt(s, heading, 0.55, y,      5.35, 0.35, size=12.5, bold=True, color=MED_GREEN)
    txt(s, detail,  0.55, y+0.38, 5.35, 0.5,  size=11,   color=DARK_TEXT)

# App screenshot — shadow rect first then image
rect(s, 6.18, 1.06, 6.8, 5.82, LIGHT_GRAY)
img(s,  P("image_1.png"), 6.1, 1.0, 6.8, 5.82)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — REAL DATA, REAL INSIGHTS
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, WHITE)
rect(s, 0, 0, 13.33, 1.0,  DARK_GREEN)

txt(s, "Real Data. Real Insights. Real Decisions.",
    0.35, 0.14, 12.6, 0.73, size=24, bold=True, color=WHITE)

# Left panel — Feed
rect(s, 0.3, 1.0, 0.07, 5.5, AMBER)
txt(s, "Feed Intelligence", 0.52, 1.05, 5.9, 0.42, size=15, bold=True, color=DARK_GREEN)
txt(s, ("Daily chart shows when birds are fed 30% below requirement — "
        "before it silently hurts production."),
    0.52, 1.52, 5.9, 0.7, size=11.5, color=DARK_TEXT)
rect(s, 0.4, 2.22, 6.05, 4.5, LIGHT_GRAY)
img(s, P("image_5.png"), 0.4, 2.22, 6.05, 4.5)

# Right panel — Finance
rect(s, 6.95, 1.0, 0.07, 5.5, MED_GREEN)
txt(s, "Financial Visibility", 7.18, 1.05, 5.8, 0.42, size=15, bold=True, color=DARK_GREEN)
txt(s, ("Top expense categories, monthly margin, all-time P&L — "
        "not at year-end but every single day."),
    7.18, 1.52, 5.8, 0.7, size=11.5, color=DARK_TEXT)
rect(s, 6.95, 2.22, 6.05, 4.5, LIGHT_GRAY)
img(s, P("image_7.png"), 6.95, 2.22, 6.05, 4.5)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BROODTRACK: WEIGHT IS THE BUSINESS
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT_BG2)
rect(s, 0, 0, 13.33, 1.0, PURPLE_DARK)

txt(s, "BroodTrack: POL/POC Farming is a Weight Business",
    0.35, 0.13, 12.6, 0.74, size=23, bold=True, color=WHITE)

# ── Left narrative panel ──────────────────────────────────────────────
rect(s, 0.3, 1.1, 5.05, 6.2, WHITE)
rect(s, 0.3, 1.1, 0.07, 6.2, PURPLE_MED)

txt(s, "You are not just raising birds.\nYou are selling readiness.",
    0.52, 1.2, 4.72, 0.9, size=15, bold=True, color=PURPLE_DARK)

txt(s, ("A POL buyer pays full price only if your birds hit"
        " target weight at the right age."
        " Underweight birds at Week 18 = lower price or no sale."),
    0.52, 2.15, 4.72, 0.85, size=11.5, color=DARK_TEXT)

points = [
    (PURPLE_MED,  "Weight tracked weekly vs breed benchmark",
     "ISA Brown, Lohmann, Noiler, Arbor Acres — each breed has a target."
     " BroodTrack shows you actual vs target every week, colour-coded."),
    (RED,         "Vaccination — one missed dose can wipe a batch",
     "Auto-generated schedule per batch. Overdue alerts appear on the dashboard"
     " before you open the app — you can't miss them."),
    (AMBER,       "Temperature during brooding is critical",
     "Week 1 target: 33–35°C. BroodTrack flags out-of-range readings"
     " immediately, so you act before chick stress causes mortality."),
    (MED_GREEN,   "Full batch P&L — cost to profit",
     "DOC purchase auto-logged. All feed, medication, and labour tracked."
     " Know your margin per batch, not just total farm revenue."),
]
for i, (col, heading, detail) in enumerate(points):
    y = 3.1 + i * 1.02
    rect(s, 0.5, y+0.06, 0.07, 0.72, col)
    txt(s, heading, 0.72, y,       4.5, 0.33, size=12, bold=True, color=PURPLE_DARK)
    txt(s, detail,  0.72, y+0.36,  4.5, 0.6,  size=10.5, color=DARK_TEXT)

# ── Right: Growth table screenshot ───────────────────────────────────
txt(s, "Actual vs Target — Week by Week (Batch LM_Q4_25 · Isa Brown)",
    5.6, 1.1, 7.45, 0.42, size=11, bold=True, color=PURPLE_DARK)
rect(s, 5.58, 1.52, 7.45, 5.75, LIGHT_GRAY)
img(s, P("image_11.png"), 5.58, 1.52, 7.45, 5.75)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — REAL BATCH, REAL RESULTS
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, WHITE)
rect(s, 0, 0, 13.33, 1.0, PURPLE_DARK)

txt(s, "Real Batch. Real Numbers. Real Profit.",
    0.35, 0.13, 12.6, 0.74, size=26, bold=True, color=WHITE)

# Batch card screenshot (left)
txt(s, "Batch overview & vaccination status",
    0.35, 1.05, 5.9, 0.38, size=11, bold=True, color=PURPLE_DARK)
rect(s, 0.3, 1.43, 5.95, 2.35, LIGHT_GRAY)
img(s, P("image_9.png"), 0.3, 1.43, 5.95, 2.35)

# Analytics report screenshot (left-bottom)
txt(s, "Batch analytics — feed, weight, efficiency, P&L",
    0.35, 3.85, 5.9, 0.38, size=11, bold=True, color=PURPLE_DARK)
rect(s, 0.3, 4.22, 5.95, 2.95, LIGHT_GRAY)
img(s, P("image_10.png"), 0.3, 4.22, 5.95, 2.95)

# Right: key stats panel
rect(s, 6.55, 1.05, 6.45, 6.3, LIGHT_BG2)
rect(s, 6.55, 1.05, 0.07, 6.3, PURPLE_MED)

txt(s, "LM_Q4_25  ·  Isa Brown  ·  Arrived 17 Nov 2025  ·  Sold",
    6.75, 1.1, 6.1, 0.38, size=10.5, bold=True, color=PURPLE_DARK)

kpis = [
    ("4,612",    "DOC In",              PURPLE_MED),
    ("4,433",    "Birds Sold (Current)", PURPLE_MED),
    ("96.1%",    "Survival Rate",        MED_GREEN),
    ("14,894 kg","Total Feed Consumed",  DARK_TEXT),
    ("3.15",     "FCR  (Pullet phase)",  AMBER),
    ("1,108g",   "Last Recorded Weight", PURPLE_MED),
]
for i, (val, lbl, col) in enumerate(kpis):
    row = i % 3
    c   = i // 3
    x   = 6.78 + c * 3.1
    y   = 1.62 + row * 1.05
    rect(s, x, y, 2.82, 0.88, WHITE)
    rect(s, x, y, 2.82, 0.07, col)
    txt(s, val, x+0.12, y+0.12, 2.6, 0.42, size=22, bold=True, color=col)
    txt(s, lbl, x+0.12, y+0.55, 2.6, 0.28, size=10,  color=GRAY)

# Profit callout
rect(s, 6.55, 4.92, 6.45, 1.62, DARK_GREEN)
rect(s, 6.55, 4.92, 0.07, 1.62, LIGHT_GREEN)
txt(s, "Net Profit — Single Batch",
    6.78, 5.0, 6.1, 0.42, size=12, bold=True, color=LIGHT_GREEN)
txt(s, "+N9,849,600",
    6.72, 5.42, 6.2, 0.72, size=34, bold=True, color=WHITE)
txt(s, "168 days  ·  4,612 DOC  ·  96.1% survival",
    6.78, 6.18, 6.1, 0.3, size=10.5, color=PALE_GREEN)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 11 — THE FUTURE (ML & PREDICTIVE)
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, DARK_GREEN)
oval(s,  9.8, -1.5, 5.0, 5.0, MED_GREEN)
oval(s, -0.8,  5.2, 3.0, 3.0, RGBColor(0x40,0x91,0x6c))

txt(s, "This is Just the Beginning",
    0.5, 0.28, 11, 0.82, size=34, bold=True, color=WHITE)
txt(s, "Every record you log today is training data for tomorrow's intelligence.",
    0.5, 1.12, 11.5, 0.52, size=16, color=LIGHT_GREEN, italic=True)
rect(s, 0.5, 1.75, 12.3, 0.05, LIGHT_GREEN)

stages = [
    ("TODAY", "Descriptive",
     RGBColor(0x40,0x91,0x6c), MED_GREEN,
     [("What happened?", True),
      "How many eggs were collected?",
      "What was my production rate?",
      "Which cells underperformed?",
      "How much did I spend on feed?",
      "",
      "→ You already have this\n   with SiloForge."]),
    ("NEXT STEP", "Diagnostic",
     BLUE, RGBColor(0x1a,0x5f,0xa8),
     [("Why did it happen?", True),
      "Why did production dip in Wk 8?",
      "Which supplier drives better FCR?",
      "What correlates with\n  high mortality?",
      "Why does Pen 2 always\n  underperform Pen 1?",
      "",
      "→ Cross-batch pattern\n   detection & root cause."]),
    ("VISION", "Predictive",
     AMBER, RGBColor(0xb8,0x7d,0x0e),
     [("What will happen?", True),
      "Predict mortality spikes\n  3 days before they occur",
      "Forecast production by\n  flock age and breed",
      "Auto-optimise feed schedules",
      "Flag disease risk from temp\n  + water + droppings data",
      "",
      "→ AI built on YOUR farm's\n   own data."]),
]
for i, (tag, title, accent, hdr_col, lines) in enumerate(stages):
    x = 0.4 + i * 4.28
    rect(s, x, 1.92, 3.98, 5.2, MED_GREEN)
    rect(s, x, 1.92, 3.98, 0.66, hdr_col)
    txt(s, tag,   x+0.15, 1.94, 2.0, 0.3,  size=10, bold=True, color=WHITE)
    txt(s, title, x+0.15, 2.26, 3.75, 0.42, size=19, bold=True, color=WHITE)
    for j, line in enumerate(lines):
        if isinstance(line, tuple):
            txt(s, line[0], x+0.2, 2.85+j*0.5, 3.6, 0.44,
                size=12.5, bold=line[1], color=accent)
        else:
            txt(s, line,    x+0.2, 2.85+j*0.5, 3.6, 0.44,
                size=11, color=PALE_GREEN)

txt(s, "→", 4.38, 4.3, 0.5, 0.5, size=26, bold=True, color=LIGHT_GREEN)
txt(s, "→", 8.66, 4.3, 0.5, 0.5, size=26, bold=True, color=LIGHT_GREEN)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BETA CTA
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, DARK_GREEN)
oval(s,  9.9, -1.3, 5.0, 5.0, MED_GREEN)
oval(s, -0.6,  5.2, 3.0, 3.0, RGBColor(0x40,0x91,0x6c))

txt(s, "Join the SiloForge Beta",
    0.5, 0.32, 10.5, 0.88, size=38, bold=True, color=WHITE)
txt(s, "FREE  ·  Limited Farm Partners  ·  Your Feedback Shapes What This Becomes",
    0.5, 1.24, 10.5, 0.5, size=15, color=LIGHT_GREEN, italic=True)
rect(s, 0.5, 1.86, 6.2, 0.05, LIGHT_GREEN)

# What you get
txt(s, "✅  What you get",
    0.5, 2.05, 5.7, 0.42, size=16, bold=True, color=LIGHT_GREEN)
get_items = [
    "Free access to LayerTrack and BroodTrack",
    "Farm setup done together today",
    "Direct line to request features",
    "Early access to every new release",
    "Become part of building Nigerian agri-tech",
]
for i, item in enumerate(get_items):
    txt(s, f"   {item}", 0.5, 2.58+i*0.57, 5.8, 0.52, size=13, color=WHITE)

# Divider
rect(s, 6.6, 2.0, 0.05, 4.9, LIGHT_GREEN)

# What we need
txt(s, "🤝  What we need from you",
    6.85, 2.05, 5.9, 0.42, size=16, bold=True, color=AMBER)
need_items = [
    "Use it consistently — log daily for 4+ weeks",
    "Tell us what is missing or broken",
    "Tell us what problems you still can't solve",
    "Be honest — good and bad feedback both help",
    "Refer us to other farmers if it works",
]
for i, item in enumerate(need_items):
    txt(s, f"   {item}", 6.85, 2.58+i*0.57, 5.9, 0.52, size=13, color=WHITE)

# Setup CTA banner
rect(s, 0.4, 5.22, 12.4, 1.95, MED_GREEN)
rect(s, 0.4, 5.22, 0.07, 1.95, LIGHT_GREEN)
txt(s, "🛠️  We are setting up farms right now",
    0.65, 5.34, 11.5, 0.48, size=16, bold=True, color=WHITE)
txt(s, ("Walk up after this session — we'll create your farm code, configure your pens,\n"
        "and you'll be logging your first data before you leave this room."),
    0.65, 5.86, 11.5, 0.95, size=13.5, color=PALE_GREEN)

# ── Save ────────────────────────────────────────────────────────────────
out = os.path.join(BASE, "SiloForge_Presentation.pptx")
prs.save(out)
print(f"\n✓ Saved: {out}")
print(f"  Slides: {len(prs.slides)}")
